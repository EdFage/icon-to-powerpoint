from fastapi import FastAPI, HTTPException
from fastapi.middleware.cors import CORSMiddleware
from fastapi.responses import FileResponse
from pydantic import BaseModel
from typing import List
import requests
import os
from pptx import Presentation
from pptx.util import Inches
from pptx.enum.text import PP_ALIGN
import tempfile
import uuid
from io import BytesIO
from PIL import Image

app = FastAPI(title="Icon to PowerPoint API")

app.add_middleware(
    CORSMiddleware,
    allow_origins=["*"],
    allow_credentials=True,
    allow_methods=["*"],
    allow_headers=["*"],
)

class IconRequest(BaseModel):
    icon_names: List[str]

def search_icon_url(icon_name: str) -> str:
    """Search for an icon and return its download URL"""
    try:
        # Try Iconify API first
        iconify_url = f"https://api.iconify.design/search?query={icon_name}&limit=1"
        response = requests.get(iconify_url)
        if response.status_code == 200:
            data = response.json()
            if data.get('icons') and len(data['icons']) > 0:
                icon_id = data['icons'][0]
                # Get SVG URL
                svg_url = f"https://api.iconify.design/{icon_id}.svg?height=128"
                return svg_url
        
        # Fallback: Use a placeholder icon service
        return f"https://via.placeholder.com/128/4A90E2/FFFFFF?text={icon_name[:2].upper()}"
        
    except Exception as e:
        print(f"Error searching for icon {icon_name}: {e}")
        return f"https://via.placeholder.com/128/4A90E2/FFFFFF?text={icon_name[:2].upper()}"

def download_icon(url: str) -> BytesIO:
    """Download icon from URL and return as BytesIO"""
    try:
        response = requests.get(url, timeout=10)
        response.raise_for_status()
        return BytesIO(response.content)
    except Exception as e:
        print(f"Error downloading icon from {url}: {e}")
        # Return a placeholder
        placeholder_response = requests.get("https://via.placeholder.com/128/4A90E2/FFFFFF?text=?")
        return BytesIO(placeholder_response.content)

@app.get("/")
async def root():
    return {"message": "Icon to PowerPoint API"}

@app.post("/search-icons")
async def search_icons(request: IconRequest):
    """Search for icons based on provided names"""
    icons_data = []
    
    for icon_name in request.icon_names:
        icon_url = search_icon_url(icon_name)
        icons_data.append({
            'name': icon_name,
            'url': icon_url
        })
    
    return {"icons": icons_data}

@app.post("/generate-powerpoint")
async def generate_powerpoint(request: IconRequest):
    """Generate PowerPoint with icons"""
    try:
        # Create presentation
        prs = Presentation()
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
        
        # Add title
        title_box = slide.shapes.add_textbox(Inches(0.5), Inches(0.5), Inches(9), Inches(1))
        title_frame = title_box.text_frame
        title_frame.text = "Your Icons"
        title_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
        title_frame.paragraphs[0].font.size = Inches(0.4)
        title_frame.paragraphs[0].font.bold = True
        
        # Calculate grid layout
        icons_per_row = 4
        icon_size = Inches(1.5)
        spacing = Inches(0.3)
        start_x = Inches(0.5)
        start_y = Inches(1.5)
        
        # Download and add icons
        for i, icon_name in enumerate(request.icon_names):
            row = i // icons_per_row
            col = i % icons_per_row
            
            x_pos = start_x + col * (icon_size + spacing)
            y_pos = start_y + row * (icon_size + spacing + Inches(0.5))
            
            # Get icon URL and download
            icon_url = search_icon_url(icon_name)
            
            try:
                # Download icon
                icon_data = download_icon(icon_url)
                
                # Save to temp file for adding to slide
                with tempfile.NamedTemporaryFile(delete=False, suffix='.png') as temp_icon:
                    # If it's SVG, we'll use placeholder for now
                    if icon_url.endswith('.svg'):
                        # Create a simple placeholder image
                        placeholder_response = requests.get(f"https://via.placeholder.com/128/4A90E2/FFFFFF?text={icon_name[:2].upper()}")
                        temp_icon.write(placeholder_response.content)
                    else:
                        temp_icon.write(icon_data.read())
                    temp_icon.flush()
                    
                    # Add image to slide
                    slide.shapes.add_picture(temp_icon.name, x_pos, y_pos, icon_size, icon_size)
                    
                    # Clean up temp file
                    os.unlink(temp_icon.name)
                    
            except Exception as e:
                print(f"Error adding icon {icon_name}: {e}")
                # Add text fallback
                text_box = slide.shapes.add_textbox(x_pos, y_pos, icon_size, icon_size)
                text_frame = text_box.text_frame
                text_frame.text = icon_name[:4].upper()
                text_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            
            # Add label below icon
            label_box = slide.shapes.add_textbox(x_pos, y_pos + icon_size + Inches(0.1), icon_size, Inches(0.3))
            label_frame = label_box.text_frame
            label_frame.text = icon_name
            label_frame.paragraphs[0].alignment = PP_ALIGN.CENTER
            label_frame.paragraphs[0].font.size = Inches(0.15)
        
        # Save to temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.pptx')
        prs.save(temp_file.name)
        temp_file.close()
        
        return FileResponse(
            temp_file.name,
            media_type='application/vnd.openxmlformats-officedocument.presentationml.presentation',
            filename='icons.pptx'
        )
        
    except Exception as e:
        raise HTTPException(status_code=500, detail=f"Error generating PowerPoint: {str(e)}")

if __name__ == "__main__":
    import uvicorn
    uvicorn.run(app, host="0.0.0.0", port=8000)